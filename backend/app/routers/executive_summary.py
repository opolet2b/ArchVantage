from fastapi import APIRouter, HTTPException
from pydantic import BaseModel
from typing import List

from app.agents.executive_summary_agent import build_executive_summary_graph

router = APIRouter(
    prefix="/executive_summary",
    tags=["executive_summary"],
)

class ExecutiveSummaryRequest(BaseModel):
    thing_id: str = None
    source_docs: List[str]
    source_asset_ids: List[str] = []
    llm_preset: str = "default"
    vlm_preset: str = "default"

import threading
import queue
import json
import asyncio
from app.core.database import SessionLocal
from app.models.canvas_models import CanvasThing
from sqlalchemy.orm.attributes import flag_modified

import logging

logger = logging.getLogger("executive_summary")
logger.setLevel(logging.INFO)
if not logger.handlers:
    ch = logging.StreamHandler()
    formatter = logging.Formatter('%(asctime)s - %(name)s - %(levelname)s - %(message)s')
    ch.setFormatter(formatter)
    logger.addHandler(ch)

@router.post("/generate")
async def run_executive_summary(request: ExecutiveSummaryRequest):
    logger.info(f"=== Incoming POST /generate request for thing_id: {request.thing_id} ===")
    
    if request.thing_id:
        try:
            with SessionLocal() as db:
                thing = db.query(CanvasThing).filter(CanvasThing.id == request.thing_id).first()
                if thing:
                    if thing.content is None:
                        thing.content = {}
                    thing.content["status"] = "generating"
                    flag_modified(thing, "content")
                    db.commit()
                    logger.info(f"Successfully locked DB status to 'generating' for thing {request.thing_id}")
                else:
                    logger.warning(f"Could not find CanvasThing with id {request.thing_id} to lock status!")
        except Exception as db_err:
            logger.error(f"Failed to set generating state to DB: {db_err}")

    q = queue.Queue()
    
    def progress_callback(completed, total):
        q.put({"type": "chunk_progress", "completed": completed, "total": total})
        
    def worker():
        logger.info(f"Background worker thread started for thing_id: {request.thing_id}")
        try:
            graph = build_executive_summary_graph()
            
            final_slides = []
            final_concepts = {}
            
            # Start streaming events from the graph
            for event in graph.stream({
                "source_docs": request.source_docs,
                "source_asset_ids": request.source_asset_ids,
                "llm_preset": request.llm_preset,
                "vlm_preset": request.vlm_preset,
                "concepts": {},
                "slides": [],
                "errors": [],
                "progress_callback": progress_callback
            }):
                # Emit step progress based on the node name
                node_name = list(event.keys())[0] if event else "unknown"
                q.put({"type": "step", "node": node_name})
                logger.info(f"Graph reached node: {node_name} for thing_id: {request.thing_id}")
                
                # If we've hit the final step, grab the slides
                if "storyboarder" in event:
                    logger.info(f"Storyboarder finished for thing_id: {request.thing_id}. Slides generated.")
                    final_slides = event["storyboarder"].get("slides", [])
                    final_concepts = event["storyboarder"].get("concepts", {})
                    q.put({"type": "completed", "slides": final_slides, "concepts": final_concepts})
            
            q.put({"type": "done"})
            logger.info(f"Graph completely finished for thing_id: {request.thing_id}")
            
            # Save results to the database if a thing_id was provided
            if request.thing_id:
                try:
                    with SessionLocal() as db:
                        thing = db.query(CanvasThing).filter(CanvasThing.id == request.thing_id).first()
                        if thing:
                            if thing.content is None:
                                thing.content = {}
                            thing.content["status"] = "completed"
                            thing.content["slides"] = final_slides
                            thing.content["concepts"] = final_concepts
                            flag_modified(thing, "content")
                            db.commit()
                            logger.info(f"Successfully saved final slides and 'completed' status to DB for thing {request.thing_id}")
                        else:
                            logger.error(f"Failed to find CanvasThing {request.thing_id} to save final results!")
                except Exception as db_err:
                    logger.error(f"Failed to save Executive Summary to DB: {db_err}")
                    
        except Exception as e:
            import traceback
            logger.error(f"Background worker CRASHED for thing_id: {request.thing_id}. Exception: {e}")
            logger.error(traceback.format_exc())
            
            q.put({"type": "error", "message": str(e)})
            
            if request.thing_id:
                try:
                    with SessionLocal() as db:
                        thing = db.query(CanvasThing).filter(CanvasThing.id == request.thing_id).first()
                        if thing:
                            if thing.content is None:
                                thing.content = {}
                            thing.content["status"] = "idle"
                            flag_modified(thing, "content")
                            db.commit()
                            logger.info(f"Successfully reset DB status to 'idle' after crash for thing {request.thing_id}")
                except Exception as db_err:
                    logger.error(f"Failed to save error state 'idle' to DB: {db_err}")

    threading.Thread(target=worker, daemon=True).start()

    async def event_stream():
        while True:
            try:
                # Use asyncio.to_thread so we don't block the event loop while waiting for queue
                msg = await asyncio.to_thread(q.get, timeout=0.1)
                yield f"data: {json.dumps(msg)}\n\n"
                if msg["type"] in ["done", "error", "completed"]:
                    break
            except queue.Empty:
                # Keep alive
                yield ": keep-alive\n\n"

    from fastapi.responses import StreamingResponse
    return StreamingResponse(event_stream(), media_type="text/event-stream")

@router.get("/status/{thing_id}")
async def get_executive_summary_status(thing_id: str):
    try:
        with SessionLocal() as db:
            thing = db.query(CanvasThing).filter(CanvasThing.id == thing_id).first()
            if not thing:
                logger.warning(f"Status check failed: Thing {thing_id} not found")
                raise HTTPException(status_code=404, detail="Thing not found")
                
            status_val = thing.content.get("status", "idle") if thing.content else "idle"
            logger.info(f"Status check for thing {thing_id}: returning status '{status_val}'")
            
            return {
                "status": status_val,
                "slides": thing.content.get("slides", []) if thing.content else [],
                "concepts": thing.content.get("concepts", {}) if thing.content else {}
            }
    except Exception as e:
        logger.error(f"Status check error: {e}")
        raise HTTPException(status_code=500, detail=str(e))

from fastapi.responses import StreamingResponse
import io

class SlideModel(BaseModel):
    title: str
    takeaway: str
    concepts: List[str]
    has_diagram: bool = False
    diagram_url: str = None
    archimate_data: dict = None

class ExportPPTXRequest(BaseModel):
    slides: List[SlideModel]

@router.post("/export_pptx")
async def export_pptx(request: ExportPPTXRequest):
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
        from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR
        import pptx
        import requests
        import io

        prs = Presentation()
        
        for slide_data in request.slides:
            # If there's a diagram, we might want a different layout, but for now we just use a blank/content layout
            slide_layout = prs.slide_layouts[1] # 1 is Title and Content
            slide = prs.slides.add_slide(slide_layout)
            
            # Title
            title_shape = slide.shapes.title
            if title_shape:
                title_shape.text = slide_data.title
                # Auto-shrink title if it's very long
                title_tf = title_shape.text_frame
                for para in title_tf.paragraphs:
                    if len(slide_data.title) > 60:
                        para.font.size = Pt(28)
                    elif len(slide_data.title) > 40:
                        para.font.size = Pt(32)
                    else:
                        para.font.size = Pt(40)
            
            # Body (Takeaway + Concepts)
            body_shape = slide.placeholders[1]
            
            # Push body down if title is exceptionally long
            top_margin = 1.5
            if len(slide_data.title) > 60:
                top_margin = 1.8
                
            body_shape.top = Inches(top_margin)
            body_shape.left = Inches(0.5)
            body_shape.width = Inches(9.0) # Default full width
            body_shape.height = Inches(7.0 - top_margin)
            
            tf = body_shape.text_frame
            tf.vertical_anchor = MSO_ANCHOR.TOP
            
            # Add Takeaway as first paragraph
            tf.text = f"Takeaway: {slide_data.takeaway}"
            p_takeaway_para = tf.paragraphs[0]
            p_takeaway_para.font.bold = True
            p_takeaway_para.font.size = Pt(16)
            p_takeaway_para.font.color.rgb = pptx.dml.color.RGBColor(0, 102, 204) # Blueish
            
            # Add Concepts as bullet points
            for concept in slide_data.concepts:
                p = tf.add_paragraph()
                p.text = concept
                p.level = 1
                p.font.size = Pt(14)
                
            tf.word_wrap = True
                
            # If there's an ArchiMate Diagram, draw it using native shapes
            if getattr(slide_data, "archimate_data", None):
                archimate = slide_data.archimate_data
                diagrams = archimate.get("diagrams", [])
                if diagrams:
                    diag = diagrams[0]
                    nodes = diag.get("nodes", [])
                    
                    # Calculate bounding box to determine scale
                    max_x = max([n.get("bounds", {}).get("x", 0) + n.get("bounds", {}).get("width", 140) for n in nodes] + [100])
                    max_y = max([n.get("bounds", {}).get("y", 0) + n.get("bounds", {}).get("height", 55) for n in nodes] + [100])
                    
                    target_w = 5.5 # Give diagrams more room
                    target_h = 5.5 
                    
                    # Dynamically calculate scale so it fits the target box
                    scale = min(target_w / max_x, target_h / max_y)
                    
                    # Base offsets
                    left_offset = Inches(4.2) # Shift further left
                    top_offset = Inches(1.5)
                    
                    shape_map = {}
                    
                    for node in nodes:
                        bounds = node.get("bounds", {"x":0,"y":0,"width":140,"height":55})
                        x = left_offset + Inches(bounds["x"] * scale)
                        y = top_offset + Inches(bounds["y"] * scale)
                        w = Inches(bounds["width"] * scale)
                        h = Inches(bounds["height"] * scale)
                        
                        node_id = node.get("id")
                        name = node_id
                        if node_id in archimate.get("elements", {}):
                            name = archimate["elements"][node_id].get("name", name)
                            
                        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
                        
                        # Style the shape to look like ArchiMate
                        shape.fill.solid()
                        shape.fill.fore_color.rgb = pptx.dml.color.RGBColor(240, 248, 255) # Light blue
                        shape.line.color.rgb = pptx.dml.color.RGBColor(100, 100, 100)
                        shape.line.width = Pt(1.0)
                        
                        # Style the text
                        tf = shape.text_frame
                        tf.text = name
                        tf.word_wrap = True
                        tf.margin_left = Pt(2)
                        tf.margin_right = Pt(2)
                        tf.margin_top = Pt(2)
                        tf.margin_bottom = Pt(2)
                        
                        # Calculate a smart font size based on the absolute height of the box
                        # h is in inches. 1 inch = 72 points. A box that is 0.5 inches tall can hold about 12pt text.
                        # h_pts = h / 914400 * 72
                        # target_font_size = min(12, max(6, h_pts / 3))
                        target_font_size = min(11, max(4, int((h / 914400) * 30)))
                        
                        # Make text fit and center
                        for para in tf.paragraphs:
                            para.alignment = PP_ALIGN.CENTER
                            para.font.size = Pt(target_font_size)
                            para.font.color.rgb = pptx.dml.color.RGBColor(30, 30, 30)
                            para.font.name = "Arial"
                            
                        shape_map[node_id] = shape
                        
                    for edge in diag.get("edges", []):
                        source_id = edge.get("source")
                        target_id = edge.get("target")
                        if source_id in shape_map and target_id in shape_map:
                            s_shape = shape_map[source_id]
                            t_shape = shape_map[target_id]
                            connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, 0, 0, 0, 0)
                            
                            # Style connector to be light/dashed to avoid visual clutter
                            connector.line.color.rgb = pptx.dml.color.RGBColor(180, 180, 180)
                            connector.line.width = Pt(1.5)
                            connector.line.dash_style = pptx.enum.dml.MSO_LINE.DASH
                            
                            # Determine best connection points based on relative positions
                            s_x = s_shape.left
                            s_y = s_shape.top
                            t_x = t_shape.left
                            t_y = t_shape.top
                            
                            if s_y < t_y - Inches(0.5): # source is above
                                connector.begin_connect(s_shape, 2) # bottom
                                connector.end_connect(t_shape, 0)   # top
                            elif s_y > t_y + Inches(0.5): # source is below
                                connector.begin_connect(s_shape, 0) # top
                                connector.end_connect(t_shape, 2)   # bottom
                            else: # roughly same level
                                if s_x < t_x: # source is left
                                    connector.begin_connect(s_shape, 3) # right
                                    connector.end_connect(t_shape, 1)   # left
                                else:
                                    connector.begin_connect(s_shape, 1) # left
                                    connector.end_connect(t_shape, 3)   # right
                            
                    # Shrink the text box to make room
                    body_shape.width = Inches(4.0)
                    
            # Else if there's a standard image, download and insert it
            elif getattr(slide_data, "diagram_url", None):
                try:
                    url = slide_data.diagram_url
                    
                    if url.startswith("data:image"):
                        # Decode base64 data URI
                        header, encoded = url.split(",", 1)
                        import base64
                        image_stream = io.BytesIO(base64.b64decode(encoded))
                    else:
                        # Handle local assets
                        if url.startswith("/api/v1/assets/"):
                            # We would need to resolve this locally, but for now we skip or construct localhost url
                            url = "http://localhost:8000" + url
                            
                        response = requests.get(url, timeout=5)
                        if response.status_code == 200:
                            image_stream = io.BytesIO(response.content)
                        else:
                            image_stream = None

                    if image_stream:
                        # Add image to the right side of the slide
                        # standard slide is 10 inches wide, 7.5 inches tall
                        left = Inches(5.2)
                        top = Inches(1.8)
                        
                        # Use PIL to get dimensions if available, otherwise just guess
                        # Since we don't know the aspect ratio, we can insert it and then resize it
                        pic = slide.shapes.add_picture(image_stream, left, top)
                        
                        # Bound picture to 4.5" x 5.0" max
                        target_w = Inches(4.5)
                        target_h = Inches(5.0)
                        
                        aspect_ratio = pic.width / max(1, pic.height)
                        target_aspect = target_w / target_h
                        
                        if aspect_ratio > target_aspect:
                            # Constrained by width
                            pic.width = target_w
                            pic.height = int(target_w / aspect_ratio)
                        else:
                            # Constrained by height
                            pic.height = target_h
                            pic.width = int(target_h * aspect_ratio)
                        
                        # Shrink the text box to make room so it doesn't overlap
                        body_shape.width = Inches(4.5)
                except Exception as e:
                    print(f"Failed to add image to PPTX: {e}")
                
        # Save to memory
        ppt_stream = io.BytesIO()
        prs.save(ppt_stream)
        ppt_stream.seek(0)
        
        return StreamingResponse(
            ppt_stream,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={
                "Content-Disposition": 'attachment; filename="executive_summary.pptx"'
            }
        )
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))

class RegenerateSlideRequest(BaseModel):
    source_docs: List[str]
    slide: SlideModel
    llm_preset: str = "default"

@router.post("/regenerate_slide")
async def regenerate_slide(request: RegenerateSlideRequest):
    try:
        from app.agents.executive_summary_agent import Slide, invoke_with_fallback, get_llm
        from langchain_core.messages import HumanMessage
        
        llm = get_llm({"llm_preset": request.llm_preset})
        
        docs_text = "\n\n".join(request.source_docs)
        
        prompt = f"""
You are an expert Executive Architect. Your task is to rewrite and polish a single slide for an executive presentation.
Use the provided Source Context as background knowledge.

Source Context:
{docs_text[:15000]} # truncate to avoid token limits

Current Slide State:
Title: {request.slide.title}
Takeaway: {request.slide.takeaway}
Bullet Points/Concepts: {', '.join(request.slide.concepts)}

Instructions:
1. Refine the Title so it is punchy and professional.
2. Rewrite the Executive Takeaway to be a compelling, 1-2 sentence summary of the slide's core message.
3. Polish the Bullet Points so they are concise, action-oriented, and clearly explain the concepts.
4. Output the result in JSON format matching the Slide schema.
"""
        res = invoke_with_fallback(llm, Slide, [HumanMessage(content=prompt)])
        
        return {
            "status": "success",
            "slide": res.model_dump()
        }
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))
