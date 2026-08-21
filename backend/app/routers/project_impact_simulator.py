from fastapi import APIRouter, HTTPException
from fastapi.responses import StreamingResponse
from pydantic import BaseModel
from typing import List, Optional

from app.agents.project_impact_agent import build_doc_parser_graph, build_constraint_solver_graph, build_copilot_graph, VariableConstraints, ExtractedTopology
from app.core.database import SessionLocal
from app.models.canvas_models import CanvasThing

router = APIRouter(
    prefix="/project_impact_simulator",
    tags=["project_impact_simulator"],
)

class IngestRequest(BaseModel):
    document_ids: List[str]
    llm_preset: str = "default"

@router.post("/ingest")
async def run_ingestion(request: IngestRequest):
    db = SessionLocal()
    try:
        # Fetch the actual document contents from DB
        documents_content = []
        for doc_id in request.document_ids:
            thing = db.query(CanvasThing).filter(CanvasThing.id == doc_id).first()
            if thing:
                # Assuming content is either text or has a text field
                content_dict = thing.content or {}
                if isinstance(content_dict, str):
                    text = content_dict
                else:
                    text = content_dict.get('text') or content_dict.get('content') or str(content_dict)
                documents_content.append(text)
        
        if not documents_content:
            return {"status": "error", "message": "No valid documents found."}

        graph = build_doc_parser_graph()
        
        # Invoke the graph
        result = graph.invoke({
            "documents": documents_content,
            "llm_preset": request.llm_preset,
            "topology": None,
            "errors": []
        })
        
        topology = result.get("topology")
        topology_data = topology.model_dump() if topology else None
        
        return {"status": "success", "report": topology_data}
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        db.close()

class SimulateRequest(BaseModel):
    thing_id: str
    topology: ExtractedTopology
    constraints: VariableConstraints
    llm_preset: str = "default"

@router.post("/simulate")
async def run_simulation(request: SimulateRequest):
    import threading
    import queue
    import json
    from sqlalchemy.orm.attributes import flag_modified
    
    try:
        with SessionLocal() as db:
            thing = db.query(CanvasThing).filter(CanvasThing.id == request.thing_id).first()
            if thing:
                if thing.content is None:
                    thing.content = {}
                thing.content["simState"] = {"step": "SIMULATING"}
                flag_modified(thing, "content")
                db.commit()
    except Exception as db_err:
        raise HTTPException(status_code=500, detail="Database lock failed")

    q = queue.Queue()
    
    def worker():
        async def process_simulation():
            import asyncio
            graph = build_constraint_solver_graph()
            
            tasks = [
                graph.ainvoke({
                    "topology": request.topology,
                    "constraints": request.constraints,
                    "llm_preset": request.llm_preset,
                    "simulation_result": None,
                    "errors": []
                })
                for _ in range(3)
            ]
            
            # Since we run 3, we can emit progress for each completion
            q.put({"type": "step", "node": "Monte Carlo Simulation (3 runs)..."})
            
            results = await asyncio.gather(*tasks, return_exceptions=True)
            
            valid_results = [r.get("simulation_result") for r in results if isinstance(r, dict) and r.get("simulation_result")]
            
            if not valid_results:
                raise Exception("All Monte Carlo simulation runs failed.")
                
            core_result = valid_results[0].model_dump()
            all_weeks = []
            all_costs = []
            WEEKLY_RATE_PER_STAFF = 3500
            
            for r in valid_results:
                if r.schedule:
                    start_w = min((c.start_week for c in r.schedule), default=0)
                    end_w = max((c.end_week for c in r.schedule), default=0)
                    calculated_weeks = max(end_w - start_w, 1)
                    calculated_cost = sum((c.end_week - c.start_week) * c.assigned_staff for c in r.schedule) * WEEKLY_RATE_PER_STAFF
                    r.total_weeks = calculated_weeks
                    r.total_cost = calculated_cost
                
                all_weeks.append(r.total_weeks)
                all_costs.append(r.total_cost)
                
            core_result["total_weeks"] = valid_results[0].total_weeks
            core_result["total_cost"] = valid_results[0].total_cost
            
            min_w = min(all_weeks)
            max_w = max(all_weeks)
            if min_w == max_w:
                min_w = int(min_w * 0.85)
                max_w = int(max_w * 1.15)
                
            min_c = min(all_costs)
            max_c = max(all_costs)
            if min_c == max_c:
                min_c = min_c * 0.85
                max_c = max_c * 1.15
            
            core_result["min_weeks_confidence"] = min_w
            core_result["max_weeks_confidence"] = max_w
            core_result["min_cost_confidence"] = min_c
            core_result["max_cost_confidence"] = max_c
            
            return core_result

        try:
            import asyncio
            loop = asyncio.new_event_loop()
            asyncio.set_event_loop(loop)
            sim_result = loop.run_until_complete(process_simulation())
            
            q.put({"type": "completed", "result": sim_result})
            
            try:
                with SessionLocal() as db:
                    thing = db.query(CanvasThing).filter(CanvasThing.id == request.thing_id).first()
                    if thing:
                        if thing.content is None:
                            thing.content = {}
                        thing.content["simState"] = {"step": "DONE"}
                        flag_modified(thing, "content")
                        db.commit()
            except Exception as db_err:
                pass
                
        except Exception as e:
            q.put({"type": "error", "message": str(e)})
            try:
                with SessionLocal() as db:
                    thing = db.query(CanvasThing).filter(CanvasThing.id == request.thing_id).first()
                    if thing:
                        if thing.content is None:
                            thing.content = {}
                        thing.content["simState"] = {"step": "WAITING"}
                        flag_modified(thing, "content")
                        db.commit()
            except Exception as db_err:
                pass
                
        q.put({"type": "done"})

    threading.Thread(target=worker, daemon=True).start()

    async def event_stream():
        import asyncio
        while True:
            try:
                msg = await asyncio.to_thread(q.get, timeout=0.1)
                yield f"data: {json.dumps(msg)}\n\n"
                if msg["type"] in ["done", "error", "completed"]:
                    break
            except queue.Empty:
                yield ": keep-alive\n\n"

    return StreamingResponse(event_stream(), media_type="text/event-stream")

@router.get("/status/{thing_id}")
async def get_sim_status(thing_id: str):
    try:
        with SessionLocal() as db:
            thing = db.query(CanvasThing).filter(CanvasThing.id == thing_id).first()
            if not thing:
                raise HTTPException(status_code=404, detail="Thing not found")
            
            content = thing.content or {}
            simState = content.get("simState", {})
            step = simState.get("step", "WAITING")
            
            return {"step": step}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

class CopilotRequest(BaseModel):
    user_message: str
    current_constraints: VariableConstraints
    llm_preset: str = "default"

@router.post("/copilot")
async def run_copilot(request: CopilotRequest):
    try:
        graph = build_copilot_graph()
        
        result = graph.invoke({
            "user_message": request.user_message,
            "current_constraints": request.current_constraints,
            "llm_preset": request.llm_preset,
            "copilot_response": None,
            "errors": []
        })
        
        copilot_response = result.get("copilot_response")
        response_data = copilot_response.model_dump() if copilot_response else None
        
        return {"status": "success", "result": response_data}
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))

class PptxExportRequest(BaseModel):
    scenario_name: str
    target_components: List[str]
    migration_pattern: str
    sim_delta: dict
    constraints: dict
    interface_protocols: dict
    topology_nodes: Optional[List[dict]] = None
    topology_edges: Optional[List[dict]] = None

class AutoSolveRequest(BaseModel):
    thing_id: str
    target_components: List[str]
    max_budget: float
    max_timeline_weeks: int
    max_staff: int
    llm_preset: str = "default"

@router.post("/auto_solve")
async def auto_solve_scenario(request: AutoSolveRequest):
    import threading
    import queue
    import json
    from sqlalchemy.orm.attributes import flag_modified
    
    try:
        with SessionLocal() as db:
            thing = db.query(CanvasThing).filter(CanvasThing.id == request.thing_id).first()
            if not thing:
                raise HTTPException(status_code=404, detail="Thing not found")
                
            if thing.content is None:
                thing.content = {}
            thing.content["simState"] = {"step": "SIMULATING"}
            flag_modified(thing, "content")
            db.commit()
            
            if not thing.content.get("report"):
                raise HTTPException(status_code=404, detail="Topology not found in CanvasThing.")
                
            topology_data = thing.content["report"]
    except Exception as db_err:
        raise HTTPException(status_code=500, detail="Database lock failed")

    q = queue.Queue()
    
    def worker():
        async def process_auto_solve():
            import asyncio
            from app.agents.project_impact_agent import ExtractedTopology, build_auto_solve_graph, AutoSolveState
            
            q.put({"type": "step", "node": "Initializing topology..."})
            topology = ExtractedTopology(**topology_data)
            
            graph = build_auto_solve_graph()
            initial_state: AutoSolveState = {
                "topology": topology,
                "target_components": request.target_components,
                "max_budget": request.max_budget,
                "max_timeline_weeks": request.max_timeline_weeks,
                "max_staff": request.max_staff,
                "llm_preset": request.llm_preset,
                "auto_solve_response": None,
                "errors": []
            }
            
            q.put({"type": "step", "node": "Solving optimization constraints (this may take up to 30s)..."})
            
            result = await asyncio.to_thread(graph.invoke, initial_state)
            
            if result.get("errors"):
                raise Exception(" ".join(result["errors"]))
                
            auto_response = result.get("auto_solve_response")
            if not auto_response:
                raise Exception("Failed to generate optimal solution.")
                
            return auto_response.model_dump()

        try:
            import asyncio
            loop = asyncio.new_event_loop()
            asyncio.set_event_loop(loop)
            solve_result = loop.run_until_complete(process_auto_solve())
            
            q.put({"type": "completed", "result": solve_result})
            
            try:
                with SessionLocal() as db:
                    thing = db.query(CanvasThing).filter(CanvasThing.id == request.thing_id).first()
                    if thing:
                        if thing.content is None:
                            thing.content = {}
                        thing.content["simState"] = {"step": "DONE"}
                        flag_modified(thing, "content")
                        db.commit()
            except Exception as db_err:
                pass
                
        except Exception as e:
            q.put({"type": "error", "message": str(e)})
            try:
                with SessionLocal() as db:
                    thing = db.query(CanvasThing).filter(CanvasThing.id == request.thing_id).first()
                    if thing:
                        if thing.content is None:
                            thing.content = {}
                        thing.content["simState"] = {"step": "WAITING"}
                        flag_modified(thing, "content")
                        db.commit()
            except Exception as db_err:
                pass
                
        q.put({"type": "done"})

    threading.Thread(target=worker, daemon=True).start()

    async def event_stream():
        import asyncio
        while True:
            try:
                msg = await asyncio.to_thread(q.get, timeout=0.1)
                yield f"data: {json.dumps(msg)}\n\n"
                if msg["type"] in ["done", "error", "completed"]:
                    break
            except queue.Empty:
                yield ": keep-alive\n\n"

    return StreamingResponse(event_stream(), media_type="text/event-stream")

@router.post("/export/pptx")
async def export_pptx(request: PptxExportRequest):
    import io
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    
    prs = Presentation()
    
    # ----------------------------------------------------------------
    # Slide 1: Title
    # ----------------------------------------------------------------
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Architecture Modernization Pitch"
    subtitle.text = f"Scenario: {request.scenario_name}\nGenerated by ArchVantage Simulator"
    
    # ----------------------------------------------------------------
    # Helper: Draw Impact Slide (Dashboard Style)
    # ----------------------------------------------------------------
    def draw_impact_slide(prs, impact_data, title_text):
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title only
        title_shape = slide.shapes.title
        title_shape.text = f"{title_text} Impact"
        
        # Draw 4 cards: Time, Budget, Risk, Bottleneck
        card_y = Inches(1.5)
        card_w = Inches(2.1)
        card_h = Inches(1.2)
        spacing = Inches(0.2)
        
        # 1. Time
        s = slide.shapes.add_shape(1, Inches(0.5), card_y, card_w, card_h)
        s.fill.solid()
        s.fill.fore_color.rgb = RGBColor(241, 245, 249) # slate-100
        s.line.color.rgb = RGBColor(203, 213, 225)
        tf = s.text_frame
        tf.word_wrap = True
        p = tf.add_paragraph()
        p.text = "TOTAL TIME"
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = RGBColor(100, 116, 139)
        p = tf.add_paragraph()
        if impact_data.get('min_weeks') and impact_data.get('max_weeks'):
            p.text = f"{impact_data.get('min_weeks_confidence', impact_data.get('min_weeks'))} - {impact_data.get('max_weeks_confidence', impact_data.get('max_weeks'))} wks"
        else:
            p.text = f"{impact_data.get('weeks')} wks"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(15, 23, 42)
        
        # 2. Budget
        s = slide.shapes.add_shape(1, Inches(0.5) + card_w + spacing, card_y, card_w, card_h)
        s.fill.solid()
        s.fill.fore_color.rgb = RGBColor(241, 245, 249)
        s.line.color.rgb = RGBColor(203, 213, 225)
        tf = s.text_frame
        p = tf.add_paragraph()
        p.text = "BUDGET"
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = RGBColor(100, 116, 139)
        p = tf.add_paragraph()
        cost_str = f"${impact_data.get('cost', 0):,}"
        if impact_data.get('min_cost') and impact_data.get('max_cost'):
            min_c = impact_data.get('min_cost_confidence', impact_data.get('min_cost'))
            max_c = impact_data.get('max_cost_confidence', impact_data.get('max_cost'))
            cost_str = f"${int(min_c/1000)}k - ${int(max_c/1000)}k" if max_c > 1000 else f"${min_c} - ${max_c}"
        p.text = cost_str
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(15, 23, 42)
        
        # 3. Risk
        s = slide.shapes.add_shape(1, Inches(0.5) + (card_w + spacing)*2, card_y, card_w, card_h)
        s.fill.solid()
        s.fill.fore_color.rgb = RGBColor(241, 245, 249)
        s.line.color.rgb = RGBColor(203, 213, 225)
        tf = s.text_frame
        p = tf.add_paragraph()
        p.text = "RISK INDEX"
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = RGBColor(100, 116, 139)
        p = tf.add_paragraph()
        p.text = str(impact_data.get('risk', 'N/A'))
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(15, 23, 42)

        # 4. Bottleneck
        s = slide.shapes.add_shape(1, Inches(0.5) + (card_w + spacing)*3, card_y, card_w, card_h)
        s.fill.solid()
        s.fill.fore_color.rgb = RGBColor(255, 241, 242) # rose-50
        s.line.color.rgb = RGBColor(254, 205, 211)
        tf = s.text_frame
        tf.word_wrap = True
        p = tf.add_paragraph()
        p.text = "CRITICAL BOTTLENECK"
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = RGBColor(225, 29, 72)
        p = tf.add_paragraph()
        bneck = str(impact_data.get('bottleneck', 'None'))
        p.text = bneck[:40] + "..." if len(bneck) > 40 else bneck
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(190, 18, 60)
        
        # Text Justifications
        text_y = Inches(3.0)
        text_w = Inches(4.3)
        text_h = Inches(3.8)
        
        # Metric Justification
        s = slide.shapes.add_shape(1, Inches(0.5), text_y, text_w, text_h)
        s.fill.solid()
        s.fill.fore_color.rgb = RGBColor(238, 242, 255) # indigo-50
        s.line.color.rgb = RGBColor(224, 231, 255)
        tf = s.text_frame
        tf.word_wrap = True
        p = tf.add_paragraph()
        p.text = "METRIC JUSTIFICATION"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = RGBColor(49, 46, 129)
        p = tf.add_paragraph()
        p.text = str(impact_data.get('justification_of_metrics', 'No justification available.'))
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(30, 41, 59)
        
        # Bottleneck Analysis
        s = slide.shapes.add_shape(1, Inches(0.5) + text_w + Inches(0.4), text_y, text_w, text_h)
        s.fill.solid()
        s.fill.fore_color.rgb = RGBColor(238, 242, 255)
        s.line.color.rgb = RGBColor(224, 231, 255)
        tf = s.text_frame
        tf.word_wrap = True
        p = tf.add_paragraph()
        p.text = "BOTTLENECK ANALYSIS"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = RGBColor(49, 46, 129)
        p = tf.add_paragraph()
        p.text = str(impact_data.get('bottleneck_analysis', impact_data.get('bottleneck_citation', 'No analysis available.')))
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(30, 41, 59)

    # Slide 2: Executive Summary (Metrics)
    delta = request.sim_delta
    draw_impact_slide(prs, delta, "Cumulative (Overall)")

    # ----------------------------------------------------------------
    # Slide 3: Scope & Strategy
    # ----------------------------------------------------------------
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Target Scope & Strategy"
    tf = body_shape.text_frame
    
    p = tf.add_paragraph()
    p.text = f"Migration Pattern: {request.migration_pattern}"
    
    c = request.constraints
    p = tf.add_paragraph()
    p.text = "Execution Strategy Constraints:"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = f"Dual Run Required: {'Yes' if c.get('dual_run') else 'No'}"
    p.level = 1
    p = tf.add_paragraph()
    p.text = f"Canary Rollout Required: {'Yes' if c.get('canary_rollout') else 'No'}"
    p.level = 1
    p = tf.add_paragraph()
    p.text = f"Zero Downtime Target: {'Yes' if c.get('zero_downtime') else 'No'}"
    p.level = 1
    p = tf.add_paragraph()
    p.text = f"Data Backfill Required: {'Yes' if c.get('data_backfill') else 'No'}"
    p.level = 1
    p = tf.add_paragraph()
    p.text = f"Assigned Team: {c.get('team_assignee')}"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Target Components:"
    for comp in request.target_components:
        p = tf.add_paragraph()
        p.text = comp
        p.level = 1

    # ----------------------------------------------------------------
    # Slide 4: Interface Protocols Mapping
    # ----------------------------------------------------------------
    if request.interface_protocols:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title only layout
        title_shape = slide.shapes.title
        title_shape.text = "Interface Protocols Strategy"
        
        # Add table
        rows = len(request.interface_protocols) + 1
        cols = 2
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(0.8)
        
        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = table_shape.table
        
        # Set column widths
        table.columns[0].width = Inches(5)
        table.columns[1].width = Inches(3)
        
        # Header
        table.cell(0, 0).text = "Dependency (Source -> Target)"
        table.cell(0, 1).text = "Protocol"
        
        # Data
        for i, (edge, proto) in enumerate(request.interface_protocols.items()):
            table.cell(i + 1, 0).text = edge
            table.cell(i + 1, 1).text = proto

    # ----------------------------------------------------------------
    # Slide 5: Visual Gantt Chart
    # ----------------------------------------------------------------
    schedule = delta.get('schedule', [])
    if schedule:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title_shape = slide.shapes.title
        title_shape.text = "Execution Schedule (Gantt)"
        
        total_weeks = delta.get('weeks', 48)
        if not total_weeks or total_weeks == 0:
            total_weeks = 48
            
        base_y = Inches(2.0)
        row_h = Inches(0.5)
        
        # Timeline Header
        txBox = slide.shapes.add_textbox(Inches(3.2), base_y - Inches(0.5), Inches(6.3), Inches(0.4))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = f"Timeline (Total: {total_weeks} weeks)"
        p.font.size = Pt(12)
        p.font.bold = True
        
        for i, s in enumerate(schedule):
            y = base_y + (i * row_h)
            
            # Label
            txBox = slide.shapes.add_textbox(Inches(0.5), y, Inches(2.6), row_h)
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.add_paragraph()
            p.text = str(s.get('component_name', s.get('component_id', ''))[:45])
            p.font.size = Pt(10)
            
            # Bar
            start = float(s.get('start_week', 0))
            end = float(s.get('end_week', 0))
            is_bot = s.get('is_bottleneck', False)
            
            gantt_x = Inches(3.2)
            gantt_w = Inches(6.3)
            
            start_x = gantt_x + (start / total_weeks) * gantt_w
            w = ((end - start) / total_weeks) * gantt_w
            w = max(w, Inches(0.1)) # min width
            
            try:
                # 1 = MSO_SHAPE.ROUNDED_RECTANGLE
                # Push the bar down slightly to align with the default textbox padding
                shape = slide.shapes.add_shape(1, start_x, y + Inches(0.12), w, row_h - Inches(0.24))
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(225, 29, 72) if is_bot else RGBColor(79, 70, 229)
                shape.line.color.rgb = RGBColor(190, 18, 60) if is_bot else RGBColor(67, 56, 202)
            except Exception as e:
                print("Failed to draw gantt bar", e)

    # ----------------------------------------------------------------
    # Slide 6+: Isolated Impacts (Topology, Org, Strategy, P&L)
    # ----------------------------------------------------------------
    isolated = delta.get('isolated_impacts', {})
    for k, title_k in [('topology', 'Topology'), ('org', 'Organization'), ('strategy', 'Execution Strategy'), ('pnl', 'P&L / Budgets')]:
        if k in isolated and isolated[k]:
            draw_impact_slide(prs, isolated[k], f"Isolated {title_k}")
    
    # ----------------------------------------------------------------
    # Slide 7: Native Topology Map
    # ----------------------------------------------------------------
    nodes = getattr(request, 'topology_nodes', [])
    edges = getattr(request, 'topology_edges', [])
    
    if nodes and len(nodes) > 0:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title only
        title_shape = slide.shapes.title
        title_shape.text = "Target Architecture Topology"
        
        # Calculate bounding box to scale to slide
        min_x = min(n.get('x', 0) for n in nodes)
        max_x = max(n.get('x', 0) for n in nodes) + 200
        min_y = min(n.get('y', 0) for n in nodes)
        max_y = max(n.get('y', 0) for n in nodes) + 70
        
        w = max_x - min_x
        h = max_y - min_y
        
        # Slide drawing area
        draw_x = Inches(0.5)
        draw_y = Inches(1.5)
        draw_w = Inches(9.0)
        draw_h = Inches(5.5)
        
        scale_x = draw_w / w if w > 0 else 1
        scale_y = draw_h / h if h > 0 else 1
        scale = min(scale_x, scale_y)
        
        # Draw Nodes First (to get shape references for connections)
        node_shapes = {}
        for n in nodes:
            nx = draw_x + (n.get('x', 0) - min_x) * scale
            ny = draw_y + (n.get('y', 0) - min_y) * scale
            nw = 200 * scale
            nh = 70 * scale
            
            # 1 = ROUNDED_RECTANGLE
            shape = slide.shapes.add_shape(1, nx, ny, nw, nh)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
            shape.line.color.rgb = RGBColor(203, 213, 225)
            
            tf = shape.text_frame
            tf.word_wrap = True
            p = tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            p.text = str(n.get('label', ''))
            
            font_multiplier = nw / Inches(2.0)
            p.font.size = Pt(max(8, int(10 * font_multiplier)))
            p.font.bold = True
            p.font.color.rgb = RGBColor(51, 65, 85)
            
            sub = n.get('subtitle')
            if sub:
                p = tf.add_paragraph()
                p.alignment = PP_ALIGN.CENTER
                p.text = f"⚡ {sub.upper()}"
                p.font.size = Pt(max(6, int(8 * font_multiplier)))
                p.font.bold = True
                p.font.color.rgb = RGBColor(245, 158, 11)
            
            node_shapes[n.get('id')] = shape

        # Draw Connectors and bind them
        edges = getattr(request, 'topology_edges', []) or []
        for e in edges:
            src_id = e.get('source')
            tgt_id = e.get('target')
            src_shape = node_shapes.get(src_id)
            tgt_shape = node_shapes.get(tgt_id)
            
            if src_shape and tgt_shape:
                # We can place them anywhere initially; begin_connect will snap them
                connector = slide.shapes.add_connector(1, 0, 0, 100, 100) # 1 = straight
                connector.line.color.rgb = RGBColor(71, 85, 105)
                connector.line.width = Pt(1.5)
                
                # Connection sites: 0=Top, 1=Left/Right, 2=Bottom, 3=Left/Right.
                # Usually 2 is bottom, 0 is top.
                try:
                    connector.begin_connect(src_shape, 2)
                    connector.end_connect(tgt_shape, 0)
                except Exception as ex:
                    print("Could not connect shapes natively:", ex)

    # ----------------------------------------------------------------
    # Save & Return
    # ----------------------------------------------------------------
    pptx_io = io.BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    
    return StreamingResponse(
        pptx_io, 
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f"attachment; filename=archvantage_{request.scenario_name.replace(' ', '_')}.pptx"}
    )
