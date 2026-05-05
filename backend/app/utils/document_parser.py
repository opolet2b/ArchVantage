import os
import io

class DocumentParser:
    """
    Utility to extract text from various file formats.
    Supported formats: .txt, .md, .csv, .pdf, .docx, .pptx
    """
    
    @staticmethod
    def extract_text_from_file(filepath: str, char_limit: int = 15000) -> str:
        """
        Extracts limited amount of raw text from a given file.
        Returns empty string if file format is unsupported or an error occurs.
        """
        if not os.path.isfile(filepath):
            return ""
            
        _, ext = os.path.splitext(filepath)
        ext = ext.lower()
        
        try:
            if ext in [".txt", ".md", ".csv", ".xml"]:
                return DocumentParser._extract_from_text(filepath, char_limit)
            elif ext in [".html", ".htm"]:
                return DocumentParser._extract_from_html(filepath, char_limit)
            elif ext == ".pdf":
                return DocumentParser._extract_from_pdf(filepath, char_limit)
            elif ext == ".docx":
                return DocumentParser._extract_from_docx(filepath, char_limit)
            elif ext == ".pptx":
                return DocumentParser._extract_from_pptx(filepath, char_limit)
            elif ext == ".xlsx":
                return DocumentParser._extract_from_xlsx(filepath, char_limit)
            elif ext in [".png", ".jpg", ".jpeg"]:
                print(f"[DocumentParser] Skipping image file {filepath} (requires multimodal/OCR)")
                return ""
            else:
                print(f"[DocumentParser] Unsupported file extension {ext} for {filepath}")
                return ""
        except Exception as e:
            print(f"[DocumentParser] Error parsing {filepath}: {e}")
            return ""

    @staticmethod
    def _extract_from_text(filepath: str, char_limit: int) -> str:
        with open(filepath, "r", encoding="utf-8", errors="replace") as f:
            return f.read(char_limit)
            
    @staticmethod
    def _extract_from_pdf(filepath: str, char_limit: int) -> str:
        try:
            from pypdf import PdfReader
            text = ""
            with open(filepath, "rb") as f:
                reader = PdfReader(f)
                for page in reader.pages:
                    text += page.extract_text() + "\n"
                    if len(text) > char_limit:
                        break
            return text[:char_limit]
        except ImportError:
            print("[DocumentParser] Missing pypdf. Cannot parse PDF.")
            return ""

    @staticmethod
    def _extract_from_docx(filepath: str, char_limit: int) -> str:
        try:
            import docx
            text = ""
            doc = docx.Document(filepath)
            for para in doc.paragraphs:
                text += para.text + "\n"
                if len(text) > char_limit:
                    break
            return text[:char_limit]
        except ImportError:
            print("[DocumentParser] Missing python-docx. Cannot parse DOCX.")
            return ""

    @staticmethod
    def _extract_from_pptx(filepath: str, char_limit: int) -> str:
        try:
            import pptx
            text = ""
            prs = pptx.Presentation(filepath)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                if len(text) > char_limit:
                    break
            return text[:char_limit]
        except ImportError:
            print("[DocumentParser] Missing python-pptx. Cannot parse PPTX.")
            return ""

    @staticmethod
    def _extract_from_xlsx(filepath: str, char_limit: int) -> str:
        try:
            import pandas as pd
            text = ""
            # Read all sheets into a dictionary of DataFrames
            excel_data = pd.read_excel(filepath, sheet_name=None)
            for sheet_name, df in excel_data.items():
                if df.empty:
                    continue
                text += f"\n### Sheet: {sheet_name}\n"
                # Convert to markdown for better LLM readability
                try:
                    text += df.to_markdown(index=False) + "\n"
                except Exception:
                    # Fallback to to_string if to_markdown fails (e.g. missing tabulate)
                    text += df.to_string(index=False) + "\n"
                
                if len(text) > char_limit:
                    break
            return text[:char_limit]
        except ImportError:
            print("[DocumentParser] Missing pandas. Cannot parse XLSX.")
            return ""
        except Exception as e:
            print(f"[DocumentParser] Excel processing error: {e}")
            return ""

    @staticmethod
    def _extract_from_html(filepath: str, char_limit: int) -> str:
        try:
            from bs4 import BeautifulSoup
            with open(filepath, "r", encoding="utf-8", errors="replace") as f:
                soup = BeautifulSoup(f, "html.parser")
                text = soup.get_text(separator="\n", strip=True)
                return text[:char_limit]
        except ImportError:
            print("[DocumentParser] Missing beautifulsoup4. Cannot parse HTML.")
            return ""
        except Exception as e:
            print(f"[DocumentParser] HTML parsing error: {e}")
            return ""

document_parser = DocumentParser()
