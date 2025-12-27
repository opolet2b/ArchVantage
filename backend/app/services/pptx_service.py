import os
from typing import List, Dict, Any, Optional
import pptx
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_AUTO_SHAPE_TYPE

class PPTXService:
    def __init__(self):
        pass

    def process_presentation(self, file_path: str) -> Dict[str, Any]:
        """
        Parse a PPTX file and return a structured JSON representation of its slides.
        
        Returns:
            Dict containing:
            - total_slides: int
            - slides: List[Dict] (each slide with index, elements, dimensions)
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"PPTX file not found: {file_path}")

        try:
            prs = pptx.Presentation(file_path)
            
            # Base dimensions (Emu units in python-pptx)
            # We want to normalize everything to 0.0 - 1.0 relative to slide size
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            
            slides_data = []
            
            for i, slide in enumerate(prs.slides):
                slide_info = self._extract_slide(slide, i, slide_width, slide_height)
                slides_data.append(slide_info)
                
            return {
                "total_slides": len(slides_data),
                "slides": slides_data,
                "meta": {
                    "width_emu": slide_width,
                    "height_emu": slide_height
                }
            }
            
        except Exception as e:
            print(f"[PPTXService] Error processing {file_path}: {e}")
            raise e

    def _extract_slide(self, slide, index: int, width_emu: int, height_emu: int) -> Dict[str, Any]:
        """Extract elements from a single slide."""
        elements = []
        
        # 1. Shapes extraction
        for shape in slide.shapes:
            element = self._parse_shape(shape, width_emu, height_emu)
            if element:
                elements.append(element)
                
        # 2. Extract speaker notes if available
        notes = ""
        if slide.has_notes_slide:
            text_frame = slide.notes_slide.notes_text_frame
            if text_frame:
                notes = text_frame.text

        return {
            "index": index,
            "elements": elements,
            "notes": notes,
            # Placeholder for AI description (Phase 2)
            "ai_description": "" 
        }

    def _extract_text_from_shape(self, shape) -> str:
        """Recursively extract text from a shape (including groups)."""
        text_parts = []
        
        # 1. Direct text frame (TextBoxes, Autoshape, Placeholders)
        if hasattr(shape, "text_frame") and shape.text_frame:
             for p in shape.text_frame.paragraphs:
                 # Flatten runs? Usually p.text is sufficient.
                 if p.text.strip():
                     text_parts.append(p.text.strip())
        
        # 2. Table text
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            try:
                for row in shape.table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text_frame and cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        text_parts.append(" | ".join(row_text))
            except Exception:
                pass

        # 3. Group recursion
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for child in shape.shapes:
                child_text = self._extract_text_from_shape(child)
                if child_text:
                    text_parts.append(child_text)
                    
        return "\n".join(text_parts)

    # Standard Office Theme Colors Fallback
    THEME_COLOR_MAP = {
        13: "FFFFFF", # LIGHT_1
        14: "000000", # DARK_1
        15: "EEECE1", # LIGHT_2
        16: "1F497D", # DARK_2
        5: "4F81BD",  # ACCENT_1 (Blue)
        6: "C0504D",  # ACCENT_2 (Red)
        7: "9BBB59",  # ACCENT_3 (Green)
        8: "8064A2",  # ACCENT_4 (Purple)
        9: "4BACC6",  # ACCENT_5 (Aqua)
        10: "F79646", # ACCENT_6 (Orange)
    }

    # Map string values from XML (e.g. 'accent1') to Hex
    STRING_THEME_MAP = {
        "lt1": "FFFFFF",
        "dk1": "000000",
        "lt2": "EEECE1",
        "dk2": "1F497D",
        "bg1": "FFFFFF",
        "tx1": "000000",
        "bg2": "EEECE1",
        "tx2": "1F497D",
        "accent1": "4F81BD",
        "accent2": "C0504D",
        "accent3": "9BBB59",
        "accent4": "8064A2",
        "accent5": "4BACC6",
        "accent6": "F79646",
    }

    def _get_color(self, shape, prop_name):
        """Helper to safely extract RGB color from fill or line."""
        try:
            # 1. Try python-pptx high-level API first
            if hasattr(shape, prop_name):
                prop = getattr(shape, prop_name)
                
                # Check for direct RGB
                try:
                    if hasattr(prop, "fore_color"):
                        if hasattr(prop.fore_color, "rgb") and prop.fore_color.rgb:
                            return str(prop.fore_color.rgb)
                        
                        # Try Theme Color (Scheme)
                        if hasattr(prop.fore_color, "type") and str(prop.fore_color.type) == "SCHEME (2)":
                             if hasattr(prop.fore_color, "theme_color"):
                                 theme_val = int(prop.fore_color.theme_color)
                                 if theme_val in self.THEME_COLOR_MAP:
                                     return self.THEME_COLOR_MAP[theme_val]
                except Exception:
                    pass

                # Check for color attribute directly (e.g. on line)
                try:
                    if hasattr(prop, "color"):
                         if hasattr(prop.color, "rgb") and prop.color.rgb:
                            return str(prop.color.rgb)
                         if hasattr(prop.color, "type") and str(prop.color.type) == "SCHEME (2)":
                             if hasattr(prop.color, "theme_color"):
                                 theme_val = int(prop.color.theme_color)
                                 if theme_val in self.THEME_COLOR_MAP:
                                     return self.THEME_COLOR_MAP[theme_val]
                except Exception:
                    pass

            # 2. Fallback: Check for Style Matrix Reference via XML (p:style/a:fillRef/a:schemeClr)
            # This handles shapes that use "Theme Style" (inheritance) and report No Fill in API
            if prop_name == "fill":
                try:
                    # p:style element exists?
                    if hasattr(shape, "_element") and hasattr(shape._element, "xpath"):
                        # Look for fillRef -> schemeClr
                        fill_refs = shape._element.xpath("p:style/a:fillRef/a:schemeClr")
                        if fill_refs:
                             val = fill_refs[0].get("val")
                             if val in self.STRING_THEME_MAP:
                                 return self.STRING_THEME_MAP[val]
                                 
                        # Also check lineRef if we were asking for line? 
                        # But this block is specifically if prop_name == "fill" (logic reused)
                except Exception:
                    pass

            if prop_name == "line":
                try:
                    if hasattr(shape, "_element") and hasattr(shape._element, "xpath"):
                        # Look for lnRef -> schemeClr
                        ln_refs = shape._element.xpath("p:style/a:lnRef/a:schemeClr")
                        if ln_refs:
                             val = ln_refs[0].get("val")
                             if val in self.STRING_THEME_MAP:
                                 return self.STRING_THEME_MAP[val]
                except Exception:
                    pass
                
            return None
        except Exception:
            return None

    def _parse_shape(self, shape, total_w, total_h) -> Optional[Dict[str, Any]]:
        """
        Parse a shape into a Normalized Element dict.
        Returns None if shape should be ignored.
        """
        try:
            # Basic geometry
            x = shape.left / total_w
            y = shape.top / total_h
            w = shape.width / total_w
            h = shape.height / total_h
            
            # For lines/connectors, these might be 0 or negative depending on direction,
            # but python-pptx normalizes .left/.top to top-left corner usually.
            # We will extract explicit start/end points for lines below.
            
            shape_type = shape.shape_type
            
            # recursive text extraction
            extracted_text = self._extract_text_from_shape(shape)
            
            # Color Extraction
            fill_color = self._get_color(shape, "fill")
            line_color = self._get_color(shape, "line")

            element = {
                "id": shape.shape_id,
                "x": x,
                "y": y,
                "w": w,
                "h": h,
                "raw_type": str(shape_type),
                "text": extracted_text,
                "fill_color": fill_color,
                "line_color": line_color,
                "rotation": getattr(shape, "rotation", 0)
            }
            
            # Identify Type & Geometry
            if shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                element["type"] = "TEXT"
                
            elif shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                 element["type"] = "TEXT" # Titles, body placeholders
                 
            elif shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                element["type"] = "SHAPE"
                if hasattr(shape, "auto_shape_type"):
                    # Convert enum value to string name if possible
                    # auto_shape_type is an int (enum member)
                    try:
                        # Find the name in MSO_AUTO_SHAPE_TYPE
                        kind_name = shape.auto_shape_type.name
                        element["shape_kind"] = kind_name
                    except Exception:
                        element["shape_kind"] = str(shape.auto_shape_type)
                        
            elif shape_type == MSO_SHAPE_TYPE.LINE:
                element["type"] = "LINE"
                # Lines have begin_x, begin_y, end_x, end_y
                # Note: these are absolute EMUs
                try:
                    element["begin_x"] = shape.begin_x / total_w
                    element["begin_y"] = shape.begin_y / total_h
                    element["end_x"] = shape.end_x / total_w
                    element["end_y"] = shape.end_y / total_h
                except Exception:
                    pass
                    
            elif shape_type == MSO_SHAPE_TYPE.PICTURE:
                element["type"] = "IMAGE"
                element["text"] = "[Image]" 
                
            elif shape_type == MSO_SHAPE_TYPE.GROUP:
                element["type"] = "GROUP"
                
            elif shape_type == MSO_SHAPE_TYPE.TABLE:
                element["type"] = "TABLE"
                
            else:
                element["type"] = "UNKNOWN"
                
            # Filter empty elements BUT keep groups if they have text
            if not extracted_text and element["type"] in ["TEXT", "TABLE", "GROUP"]:
                 return None
            
            # Keep visual shapes (SHAPE, LINE, IMAGE, UNKNOWN) even if no text
            
            return element
            
        except Exception as e:
            print(f"[PPTXService] Warning parsing shape: {e}")
            return None

pptx_service = PPTXService()
